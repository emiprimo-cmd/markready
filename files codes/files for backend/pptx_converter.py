from pptx import Presentation

def convert_pptx(filepath: str) -> str:
    """Convert PPTX to markdown, one section per slide."""
    prs = Presentation(filepath)
    lines = []

    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"## Slide {i}")

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                # Detect title placeholder
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    continue
                if shape.name.lower().startswith("title"):
                    lines.append(f"### {text}")
                else:
                    lines.append(f"- {text}")

        lines.append("")  # blank line between slides

    return "\n".join(lines)
